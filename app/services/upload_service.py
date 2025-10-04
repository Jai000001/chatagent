from fastapi import UploadFile
from typing import List

async def handle_upload_files(files: List[UploadFile], task_id, client_id, dept_ids, duplicate_file_name, url_uuid, postgres_handler):
    import time
    from app.core.logger import Logger
    logger = Logger.get_logger(__name__)
    from app.utils.file_utils_manager.file_utils import FileManager
    from app.adapters.file_processor import FileProcessor
    from app.adapters.database.redisdb_handler import RedisDBHandler
    from app.adapters.database.qdrantdb_handler import QdrantDBHandler
    from app.utils.text_splitter import shared_text_splitter
    from app.utils.shared_utils import delete_existing_source
    from app.core.app_config import app_config
    import aiofiles
    import asyncio
    import os
    from PyPDF2 import PdfReader
    from docx import Document as DocxDocument
    from pptx import Presentation

    redis_handler = RedisDBHandler()
    qdrant_handler = QdrantDBHandler()

    file_processor = FileProcessor()
    file_manager = FileManager()
    start_time = time.time()
    
    if not files:
        logger.error({"error": "No files provided"})
        raise ValueError("No files provided")

    docs = []
    tasks = []
    total_pages = 0
    processed_pages = 0
    file_page_counts = []
    failed_files = []
    saved_files = []
    upload_dir = os.path.join("/tmp", "uploads")
    os.makedirs(upload_dir, exist_ok=True)

    # Async concurrent save for UploadFile objects
    semaphore = asyncio.Semaphore(100)  # Limit concurrency for memory safety
    async def save_and_dedup(upload: UploadFile):
        async with semaphore:
            filename = upload.filename
            file_path = os.path.join(upload_dir, filename)
            try:
                async with aiofiles.open(file_path, 'wb') as out_file:
                    content = await upload.read()
                    await out_file.write(content)
                # Duplicate detection
                if await file_manager.is_duplicate_file(file_path, client_id):
                    logger.info(f"Duplicate file detected and skipped: {filename}")
                    return None
                return (filename, file_path)
            except Exception as e:
                logger.error(f"Failed to save or deduplicate file {filename}: {e}")
                failed_files.append({"filename": filename, "error": f"Failed to save or deduplicate: {str(e)}"})
                return None

    # If files are UploadFile objects, save concurrently
    if files and isinstance(files[0], UploadFile):
        results = await asyncio.gather(*(save_and_dedup(upload) for upload in files))
        saved_files = [r for r in results if r]
    # If files are already (filename, file_path) tuples (from ARQ), skip saving
    elif files and isinstance(files[0], tuple):
        saved_files = files
    elif files and isinstance(files[0], list):
        saved_files = files[0]
    else:
        logger.error(f"[ERROR] Unhandled files input type: {type(files)}, value: {files}")
        raise ValueError("Could not process uploaded files: unknown input type.")

    # Calculate total number of pages/slides across all files
    for (filename, file_path) in saved_files:
        file_extension = os.path.splitext(filename)[1].lower()
        try:
            if file_extension == '.pdf':
                pdf_reader = PdfReader(file_path)
                num_pages = len(pdf_reader.pages)
            elif file_extension == '.docx':
                doc = DocxDocument(file_path)
                num_pages = len(doc.paragraphs)
            elif file_extension == '.pptx':
                ppt = Presentation(file_path)
                num_pages = len(ppt.slides)
            elif file_extension == '.txt':
                try:
                    from langchain_community.document_loaders.helpers import detect_file_encodings
                    detected_encodings = detect_file_encodings(file_path)
                    for encoding in detected_encodings:
                        try:
                            with open(file_path, 'r', encoding=encoding.encoding) as file:
                                num_pages = len(file.readlines())
                                break
                        except UnicodeDecodeError:
                            continue
                except Exception as e:
                    with open(file_path, 'r') as file:
                        num_pages = len(file.readlines())
            elif file_extension == '.doc':
                from app.adapters.file_parsers.doc_parser import DOCParser
                doc_parser = DOCParser()
                docx_path = doc_parser._convert_doc_to_docx(file_path)
                doc = DocxDocument(docx_path)
                num_pages = len(doc.paragraphs)
                file_path = docx_path
            elif file_extension == '.ppt':
                from app.adapters.file_parsers.ppt_parser import PPTParser
                ppt_parser = PPTParser()
                pptx_path = ppt_parser._convert_ppt_to_pptx(file_path)
                ppt = Presentation(pptx_path)
                num_pages = len(ppt.slides)
                file_path = pptx_path
            else:
                logger.error(f"Unsupported file type: {file_extension}")
                failed_files.append({"filename": filename, "error": f"Unsupported file type: {file_extension}"})
                continue

            if num_pages == 0:
                logger.error(f"No content found in file {file_path}")
                failed_files.append({"filename": filename, "error": "No content found in file"})
                continue

            file_page_counts.append((filename, file_path, num_pages))
            total_pages += num_pages

        except Exception as e:
            logger.error(f"Error reading file {file_path}: {e}")
            failed_files.append({"filename": filename, "error": f"Error reading file: {str(e)}"})
            continue

    if total_pages == 0 and not failed_files:
        logger.error({"error": "No pages or slides found in files"})
        raise ValueError("No pages or slides found in files")

    # Controlled parallelism: process up to 4 files at a time
    semaphore = asyncio.Semaphore(4)

    async def process_file_with_limit(filename, file_path, num_pages):
        nonlocal processed_pages
        async with semaphore:
            if duplicate_file_name:
                duplicate_file_list = [f.strip() for f in duplicate_file_name.split(",") if f.strip()]
                if filename in duplicate_file_list:
                    dept_ids_list = [dept_ids] if isinstance(dept_ids, str) else dept_ids if isinstance(dept_ids, list) else []
                    try:
                        await delete_existing_source(client_id, filename, dept_ids_list)
                    except Exception as e:
                        logger.warning(f"Failed to delete existing source for {filename}: {e}")
            try:
                await file_processor.load_file_async(file_path, docs, filename)
                for page_num in range(num_pages):
                    processed_pages += 1
                    await file_manager.update_upload_file_progress(task_id, processed_pages, total_pages)
            except Exception as e:
                logger.error(f"Error processing file {filename}: {e}")
                failed_files.append({"filename": filename, "error": f"Error processing file: {str(e)}"})

    # Schedule up to 4 files in parallel
    for idx, (filename, file_path, num_pages) in enumerate(file_page_counts):
        if os.path.exists(file_path) and os.path.getsize(file_path) > 0:
            tasks.append(process_file_with_limit(filename, file_path, num_pages))
        else:
            logger.error(f"File {filename} is empty or could not be saved.")
            failed_files.append({"filename": filename, "error": "File is empty or could not be saved"})

    if tasks:
        try:
            await asyncio.gather(*tasks)
        except Exception as e:
            logger.error(f"Error during file processing: {e}")

    if not docs and not failed_files:
        await redis_handler.set_progress_in_store(task_id, {"progress": 100, "execution_time": None, "end_time": None, "process_type": "files", "total_tokens": 0, "total_cost": 0, "uploaded_files_details": []})
        logger.error({"error": "No files processed"})
        raise ValueError("No files processed")

    uploaded_files_details = []

    # Only process further if we have any valid documents
    if docs:
        # Store documents in PostgreSQL before further processing
        for doc in docs:
            if 'task_id' not in doc.metadata:
                doc.metadata['task_id'] = task_id
            success = await postgres_handler.store_upload_files_document(doc)
            if not success:
                logger.warning(f"Failed to store document in PostgreSQL: {doc.metadata.get('source', 'unknown')}")

        # First, split the documents into chunks
        all_splits = shared_text_splitter.split_documents(docs)

        # Deduplicate
        unique_splits = []
        duplicate_count = 0

        for split in all_splits:
            if 'task_id' not in split.metadata:
                split.metadata['task_id'] = task_id
            # Deduplication using file_manager
            if await file_manager.is_duplicate_file(split, client_id):
                duplicate_count += 1
            else:
                unique_splits.append(split)
                await postgres_handler.store_upload_files_document(split)

        def calculate_file_cost(filename, docs):
            file_tokens = sum(qdrant_handler.calculate_tokens(doc.page_content) for doc in docs if doc.metadata.get('source', '').endswith(filename))
            cost_per_1k_tokens = app_config.EMBEDDING_MODEL_RATE_PER_1K_TOKENS
            embedding_cost = (file_tokens / 1000) * cost_per_1k_tokens
            return f"{embedding_cost:.6f}", file_tokens

        def calculate_file_statistics(docs, filename):
            file_content = ' '.join(doc.page_content for doc in docs if doc.metadata.get('source', '').endswith(filename))
            char_count = len(file_content)
            word_count = len(file_content.split())
            return char_count, word_count

        async def modified_add_documents(unique_splits):
            file_details = []
            for filename, _, _ in file_page_counts:
                file_docs = [doc for doc in unique_splits if doc.metadata.get('source', '').endswith(filename)]
                if file_docs:
                    file_cost, file_tokens = calculate_file_cost(filename, file_docs)
                    char_count, word_count = calculate_file_statistics(file_docs, filename)
                    file_details.append({
                        "filename": filename,
                        "cost": file_cost,
                        "tokens": file_tokens,
                        "size_mb": None,
                        "character_count": char_count,
                        "word_count": word_count,
                        "status": "uploaded",
                        "error": None,
                        "reason": None
                    })
            # Add documents to QdrantDB
            await qdrant_handler.add_documents(unique_splits, task_id, client_id, dept_ids, url_uuid)
            for file_detail in file_details:
                try:
                    size = await qdrant_handler.get_collection_size_mb(file_detail['filename'], client_id, dept_ids)
                    file_detail['size_mb'] = f"{size:.6f}"
                except Exception as e:
                    logger.warning(f"Could not retrieve size for {file_detail['filename']}: {e}")
                    file_detail['size_mb'] = None
            return file_details

        uploaded_files_details = await modified_add_documents(unique_splits)
    
    for failed_file in failed_files:
        uploaded_files_details.append({
            "filename": failed_file["filename"],
            "cost": "0",
            "tokens": 0,
            "size_mb": "0",
            "character_count": 0,
            "word_count": 0,
            "status": "failed",
            "error": failed_file["error"],
            "reason": "File could not be read."
        })

    failed_files.clear()  # Clear the list to avoid memory issues

    end_time = time.time()
    execution_time_round = round(end_time - start_time, 2)
    execution_time = f"{execution_time_round} seconds"
    formatted_end_time = time.strftime("%H:%M:%S", time.localtime(end_time))

    logger.info(f"Total chunks processed: {len(all_splits) if docs else 0}")
    logger.info(f"Unique chunks after deduplication: {len(unique_splits) if docs else 0}")
    logger.info({"task_id": task_id, "message": "Files have been processed"})

    current_progress = await redis_handler.get_progress_from_store(task_id) or {}
    current_progress.update({
        "progress": 100,
        "execution_time": execution_time,
        "end_time": formatted_end_time,
        "total_tokens": current_progress.get("total_tokens", 0),
        "total_cost": current_progress.get("total_cost", 0),
        "uploaded_files_details": uploaded_files_details
    })
    await redis_handler.set_progress_in_store(task_id, current_progress)

    return {
        "task_id": task_id,
        "message": "Files have been processed",
        "total_chunks": len(all_splits) if docs else 0,
        "unique_chunks": len(unique_splits) if docs else 0,
        "duplicates_removed": duplicate_count if docs else 0,
        "uploaded_files_details": uploaded_files_details,
        "failed_files_count": len(failed_files)
    }
