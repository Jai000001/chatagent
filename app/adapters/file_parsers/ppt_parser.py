from app.core.logger import Logger
logger = Logger.get_logger(__name__)

class PPTParser:
    def __init__(self):
        self.psm_config = '--psm 6'
        self.min_text_length = 5
        self.max_width = 1000

    async def load_async(self, file_path):
        # try:
        import asyncio
        from langchain_community.document_loaders import UnstructuredPowerPointLoader
        loader = UnstructuredPowerPointLoader(file_path)
        docs = loader.load()
        text = "\n".join(doc.page_content for doc in docs)
        loop = None
        try:
            loop = asyncio.get_running_loop()
        except RuntimeError:
            loop = asyncio.new_event_loop()
            asyncio.set_event_loop(loop)
        image_text = await loop.run_in_executor(None, self._extract_text_from_ppt_images, file_path)
        return text + "\n" + image_text
        # except Exception as e:
        #     logger.error(f"Exception while loading ppt: {e}")
        #     return ""

    def _convert_ppt_to_pptx(self, input_file):
        try:
            import aspose.slides as slides
            output_file = input_file.rsplit('.', 1)[0] + '.pptx'
            with slides.Presentation(input_file) as presentation:
                presentation.save(output_file, slides.export.SaveFormat.PPTX)
            return output_file
        except Exception as e:
            logger.error(f"Exception while converting ppt to pptx: {e}")
            return input_file

    def _extract_text_from_ppt_images(self, file_path):
        # try:
        from pptx import Presentation
        from concurrent.futures import ThreadPoolExecutor, as_completed
        image_text = []
        prs = Presentation(file_path)
        with ThreadPoolExecutor() as executor:
            future_to_slide = {executor.submit(self._process_slide, slide, slide_num): slide_num
                            for slide_num, slide in enumerate(prs.slides, start=1)}
            for future in as_completed(future_to_slide):
                slide_num = future_to_slide[future]
                slide_text = future.result()
                if slide_text:
                    image_text.append(f"Slide {slide_num}:\n" + "\n".join(slide_text))
        return "\n\n".join(image_text)
        # except Exception as e:
        #     logger.error(f"Exception while extracting text from ppt image: {e}")
        #     return ""

    def _process_slide(self, slide, slide_num):
        slide_text = []
        try:
            import io, numpy as np, cv2, pytesseract
            for shape_num, shape in enumerate(slide.shapes, start=1):
                if hasattr(shape, 'image'):
                    try:
                        img_stream = io.BytesIO(shape.image.blob)
                        img = Image.open(img_stream)
                        img_np = np.array(img)
                        if len(img_np.shape) == 2 or (len(img_np.shape) == 3 and img_np.shape[2] == 1):
                            img_cv = img_np
                        elif len(img_np.shape) == 3:
                            if img_np.shape[2] == 3:
                                try:
                                    img_cv = cv2.cvtColor(img_np, cv2.COLOR_BGR2GRAY)
                                except cv2.error:
                                    img_cv = cv2.cvtColor(img_np, cv2.COLOR_RGB2GRAY)
                            elif img_np.shape[2] == 4:
                                img_cv = cv2.cvtColor(img_np, cv2.COLOR_RGBA2GRAY)
                            else:
                                raise ValueError(f"Unexpected number of channels: {img_np.shape[2]}")
                        else:
                            raise ValueError(f"Unexpected image shape: {img_np.shape}")
                        text = self._quick_ocr(img_cv)
                        if len(text.strip()) >= self.min_text_length:
                            slide_text.append(f"Shape {shape_num}: {text.strip()}")
                        else:
                            text = self._detailed_ocr(img_cv)
                            if len(text.strip()) >= self.min_text_length:
                                slide_text.append(f"Shape {shape_num}: {text.strip()}")
                    except Exception as e:
                        logger.error(f"Error processing image in Slide {slide_num}, Shape {shape_num}: {str(e)}")
                        continue
        except Exception as e:
            logger.error(f"Exception in _process_slide: {e}")
        return slide_text

    def _quick_ocr(self, img):
        try:
            import pytesseract
            return pytesseract.image_to_string(img, config=self.psm_config)
        except Exception as e:
            logger.error(f"Exception while quick ocr: {e}")
            return ""

    def _detailed_ocr(self, img):
        try:
            import cv2, pytesseract
            denoised = cv2.fastNlMeansDenoising(img)
            thresh = cv2.threshold(denoised, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)[1]
            return pytesseract.image_to_string(thresh, config=self.psm_config)
        except Exception as e:
            logger.error(f"Exception while detailed ocr: {e}")
            return ""
